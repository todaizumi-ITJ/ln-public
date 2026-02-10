#!/usr/bin/env python3
"""
LNポータルシステム マニュアルのPowerPoint作成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # スライドマスター
    blank_layout = prs.slide_layouts[6]  # 空白レイアウト
    
    # ===== スライド1: タイトル =====
    slide = prs.slides.add_slide(blank_layout)
    
    # 背景色（グラデーション風）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(37, 99, 235)  # ブルー
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "LNポータルシステム"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "完全操作マニュアル"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # バージョン
    version_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.4))
    version_frame = version_box.text_frame
    version_frame.text = "バージョン 1.0 | 2026-02-10"
    version_para = version_frame.paragraphs[0]
    version_para.font.size = Pt(20)
    version_para.font.color.rgb = RGBColor(200, 220, 255)
    version_para.alignment = PP_ALIGN.CENTER
    
    # ===== スライド2: システム概要 =====
    slide = prs.slides.add_slide(blank_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    # タイトル
    add_slide_title(slide, "システム概要")
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "📊 LNポータルは、発信者情報開示請求・裁判業務を効率化するシステムです", 0)
    add_bullet(tf, "主な機能", 0, bold=True)
    add_bullet(tf, "案件（LN）の一元管理", 1)
    add_bullet(tf, "申立書の自動生成", 1)
    add_bullet(tf, "業務マニュアルの提供", 1)
    add_bullet(tf, "実績・レベル管理（ゲーミフィケーション）", 1)
    add_bullet(tf, "フィードバック・改善提案", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "対象ユーザー", 0, bold=True)
    add_bullet(tf, "🛡️ beginner（初級者）: 証拠収集、LNチェック", 1)
    add_bullet(tf, "⚔️ expert（中級者）: 申立書作成、書類詰め", 1)
    add_bullet(tf, "🏰 leader（上級者）: 全体管理、発送、指導", 1)
    
    # ===== スライド3: ユーザー選択とホーム画面 =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    add_slide_title(slide, "ユーザー選択とホーム画面")
    
    # スクリーンショット
    img_path = "guide-screenshots/annotated-01-home.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5))
    
    # 説明テキスト
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "① ユーザー選択", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "画面上部で自分の名前を選択", 1)
    add_bullet(tf, "ロールに応じた画面が表示される", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "② LN一覧クリック", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "サイドバーから「📄 LN一覧」をクリック", 1)
    add_bullet(tf, "全案件の一覧が表示される", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "③ サマリー確認", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "全LN: 総案件数", 1)
    add_bullet(tf, "今日やること: 本日対応が必要な案件数", 1)
    add_bullet(tf, "完了・処理中: 進捗状況", 1)
    
    # ===== スライド4: LN一覧画面 =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    add_slide_title(slide, "LN一覧画面")
    
    # スクリーンショット
    img_path = "guide-screenshots/annotated-02-ln-list.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5))
    
    # 説明テキスト
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "① タブ切替", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "全LN: すべての案件を表示", 1)
    add_bullet(tf, "今日やること: 本日対応が必要な案件のみ", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "② 検索", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "訴番・申立人・相手方で検索", 1)
    add_bullet(tf, "リアルタイムフィルタリング", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "③ カードクリック", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "LNカードをクリックで詳細表示", 1)
    add_bullet(tf, "証拠進捗・書類状況を確認", 1)
    add_bullet(tf, "「詳細を見る」で展開", 1)
    
    # ===== スライド5: 申立書作成画面 =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    add_slide_title(slide, "申立書作成（makeln）画面")
    
    # スクリーンショット
    img_path = "guide-screenshots/annotated-03-makeln.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5))
    
    # 説明テキスト
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "① 種類選択", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "開示のみ", 1)
    add_bullet(tf, "開示＋消去禁止", 1)
    add_bullet(tf, "開示＋提供命令", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "② パターン選択", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "著作権証明パターンを選択", 1)
    add_bullet(tf, "10種類のパターンから選択", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "③ 作成実行", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "「📄 申立書を作成」をクリック", 1)
    add_bullet(tf, "PDF/Markdownでダウンロード", 1)
    
    # ===== スライド6: マニュアル画面 =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    add_slide_title(slide, "マニュアル画面")
    
    # スクリーンショット
    img_path = "guide-screenshots/annotated-04-manual.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5))
    
    # 説明テキスト
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "① タブ選択", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "はじめに / 証拠の集め方", 1)
    add_bullet(tf, "LN確認 / 書類詰め", 1)
    add_bullet(tf, "expert業務 / leader業務", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "② 項目展開", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "アコーディオン形式で展開", 1)
    add_bullet(tf, "手順や説明が表示される", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "③ フィードバック", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "各項目に「📝 改善提案」ボタン", 1)
    add_bullet(tf, "マニュアルの改善を提案できる", 1)
    
    # ===== スライド7: 実績・レベル画面 =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    
    add_slide_title(slide, "実績・レベル画面")
    
    # スクリーンショット
    img_path = "guide-screenshots/annotated-05-game.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5))
    
    # 説明テキスト
    text_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "① タブ切替", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "📊 マイページ: レベル・ポイント確認", 1)
    add_bullet(tf, "📋 日次レポート: 本日の作業報告", 1)
    add_bullet(tf, "🏆 リーダーボード: ランキング", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "② 作業入力", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "今日の作業内容を記入", 1)
    add_bullet(tf, "伝達事項を共有", 1)
    
    tf.add_paragraph()
    add_bullet(tf, "③ 送信", 0, bold=True, color=RGBColor(220, 38, 38))
    add_bullet(tf, "「送信」ボタンで保存", 1)
    add_bullet(tf, "自動的にポイント加算", 1)
    
    # ===== スライド8: まとめ =====
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(37, 99, 235)
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "まとめ"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.333), Inches(3.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "✅ システムの基本操作を理解しましょう", 0, size=28, color=RGBColor(255, 255, 255))
    add_bullet(tf, "✅ 自分のロールに応じた業務を確認しましょう", 0, size=28, color=RGBColor(255, 255, 255))
    add_bullet(tf, "✅ マニュアルを活用して効率的に作業しましょう", 0, size=28, color=RGBColor(255, 255, 255))
    add_bullet(tf, "✅ 実績・レベルでモチベーションを維持しましょう", 0, size=28, color=RGBColor(255, 255, 255))
    add_bullet(tf, "✅ フィードバックでシステムを改善していきましょう", 0, size=28, color=RGBColor(255, 255, 255))
    
    # アクセスURL
    url_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.6))
    url_frame = url_box.text_frame
    url_frame.text = "https://todaizumi-itj.github.io/ln-public/lnportal.html"
    url_para = url_frame.paragraphs[0]
    url_para.font.size = Pt(24)
    url_para.font.color.rgb = RGBColor(200, 220, 255)
    url_para.alignment = PP_ALIGN.CENTER
    
    # 保存
    output_path = "LNポータル_操作マニュアル.pptx"
    prs.save(output_path)
    print(f"✓ PowerPoint作成完了: {output_path}")
    return output_path

def add_slide_title(slide, title_text):
    """スライドにタイトルを追加"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(37, 99, 235)
    
    # 下線
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(1.4),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(37, 99, 235)
    line.line.color.rgb = RGBColor(37, 99, 235)

def add_bullet(text_frame, text, level=0, bold=False, size=20, color=None):
    """箇条書きを追加"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    if bold:
        p.font.bold = True
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(30, 41, 59)

if __name__ == "__main__":
    create_presentation()
